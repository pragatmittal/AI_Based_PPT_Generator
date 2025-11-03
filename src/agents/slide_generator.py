import os
from pptx import Presentation
from pptx.util import Inches
from dotenv import load_dotenv
class SlideGenerator:
    def __init__(self,output_dir="output/"):
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)
    def create_slide_deck(self,title,bullet_points):
        prs=Presentation()
        #Title Slide
        slide_layout=prs.slide_layouts[0]
        slide =prs.slides.add_slide(slide_layout)
        title_placeholder=slide.shapes.title
        title_placeholder.text=title

        #content Slide
        for i in range(0,len(bullet_points),5):
            slide_layout=prs.slide_layouts[1]
            slide =prs.slides.add_slide(slide_layout)
            title_shape=slide.shapes.title
            content_shape=slide.shapes.placeholders[1]
            title_shape.text=f"{title}(part {i//5 +1})"
            for bullet in bullet_points[i:i+5]:
                p=content_shape.text_frame.add_paragraph()
                p.text=bullet
            
        #save presentation
        pptx_path=os.path.join(self.output_dir,"generated_presentation.pptx")
        prs.save(pptx_path)
        print(f"\nPresentation saved at {pptx_path}")
        
if __name__=="__main__":
    sample_title="Sample Presentation"
    sample_points=[
            "Point 1: Introduction to the topic.",
            "Point 2: Key concepts and definitions.",
            "Point 3: Detailed analysis of the subject matter.",
            "Point 4: Case studies and real-world examples.",
            "Point 5: Summary of findings and conclusions.",
            "Point 6: Future directions and recommendations.",
            "Point 7: Q&A session and discussion points.",
            "Point 8: Additional resources and references.",
            "Point 9: Acknowledgments and credits.",
            "Point 10: Closing remarks and final thoughts."
        ]
    slide_generator=SlideGenerator()
    slide_generator.create_slide_deck(sample_title,sample_points)
        